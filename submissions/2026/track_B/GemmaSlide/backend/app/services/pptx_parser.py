from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.schemas.pptx import BBoxEmu, ShapeElement, SlideResult, TextParagraph
from app.services.coordinates import emu_to_norm, emu_to_px
from app.services.image_service import SlideImageInfo


def shape_type_name(shape_type_code: int) -> str:
    try:
        return MSO_SHAPE_TYPE(shape_type_code).name
    except Exception:
        return f"UNKNOWN_{shape_type_code}"


def shape_bbox_emu(shape: Any) -> BBoxEmu:
    return BBoxEmu(
        x=int(getattr(shape, "left", 0) or 0),
        y=int(getattr(shape, "top", 0) or 0),
        width=int(getattr(shape, "width", 0) or 0),
        height=int(getattr(shape, "height", 0) or 0),
    )


def extract_text_payload(shape: Any) -> tuple[str | None, list[TextParagraph]]:
    if not getattr(shape, "has_text_frame", False):
        return None, []
    if not shape.has_text_frame:
        return None, []

    paragraphs: list[TextParagraph] = []
    all_text: list[str] = []

    for para in shape.text_frame.paragraphs:
        run_texts = [run.text for run in para.runs if run.text]
        para_text = "".join(run_texts).strip() if run_texts else (para.text or "").strip()
        if para_text:
            all_text.append(para_text)
        paragraphs.append(TextParagraph(text=para_text, runs=run_texts))

    merged = "\n".join([t for t in all_text if t]).strip()
    return (merged if merged else None), paragraphs


def extract_shape_extra(shape: Any) -> dict[str, Any]:
    extra: dict[str, Any] = {}
    if getattr(shape, "is_placeholder", False):
        extra["is_placeholder"] = True

    if getattr(shape, "has_chart", False):
        extra["has_chart"] = True

    if getattr(shape, "has_table", False):
        rows = len(shape.table.rows)
        cols = len(shape.table.columns)
        extra["table_rows"] = rows
        extra["table_cols"] = cols

    return extra


def extract_image_payload(shape: Any) -> tuple[str | None, str | None, int | None]:
    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.PICTURE:
        return None, None, None

    image = getattr(shape, "image", None)
    if image is None:
        return None, None, None

    return image.content_type, image.ext, len(image.blob)


def parse_presentation(
    pptx_path: str,
    image_infos: dict[int, SlideImageInfo],
    flatten_groups: bool,
    element_types: set[str] | None,
) -> tuple[list[SlideResult], int]:
    prs = Presentation(pptx_path)
    slide_width_emu = int(prs.slide_width)
    slide_height_emu = int(prs.slide_height)

    slides: list[SlideResult] = []
    total_elements = 0

    if len(prs.slides) == 0:
        return [], 0

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_warnings: list[str] = []
        slide_elements: list[ShapeElement] = []
        image_info = image_infos.get(slide_index)

        def parse_shape(
            shape: Any,
            z_index: int,
            level: int,
            parent_id: str | None,
            path: str,
        ) -> None:
            nonlocal total_elements

            shp_type = int(getattr(shape, "shape_type", -1))
            shp_type_name = shape_type_name(shp_type)

            bbox_emu = shape_bbox_emu(shape)
            bbox_norm = emu_to_norm(bbox_emu, slide_width_emu, slide_height_emu)
            bbox_px = None
            if image_info is not None:
                bbox_px = emu_to_px(
                    bbox_emu,
                    slide_width_emu,
                    slide_height_emu,
                    image_info.width_px,
                    image_info.height_px,
                )

            text, paragraphs = extract_text_payload(shape)
            extra = extract_shape_extra(shape)
            media_type, image_ext, image_size_bytes = extract_image_payload(shape)

            element = ShapeElement(
                element_id=path,
                parent_id=parent_id,
                level=level,
                z_index=z_index,
                name=getattr(shape, "name", "unnamed"),
                shape_type_code=shp_type,
                shape_type_name=shp_type_name,
                is_group=(shp_type == MSO_SHAPE_TYPE.GROUP),
                has_text=bool(text),
                text=text,
                paragraphs=paragraphs,
                rotation=float(getattr(shape, "rotation", 0.0)),
                bbox_emu=bbox_emu,
                bbox_norm=bbox_norm,
                bbox_px=bbox_px,
                media_type=media_type,
                image_ext=image_ext,
                image_size_bytes=image_size_bytes,
                table_rows=extra.get("table_rows"),
                table_cols=extra.get("table_cols"),
                has_chart=bool(extra.get("has_chart", False)),
                extra=extra,
            )

            include = True
            if element_types:
                include = element.shape_type_name.lower() in element_types

            if include:
                slide_elements.append(element)
                total_elements += 1

            if shp_type == MSO_SHAPE_TYPE.GROUP:
                child_shapes = list(getattr(shape, "shapes", []))
                for child_idx, child in enumerate(child_shapes, start=1):
                    child_path = f"{path}.{child_idx}"
                    if flatten_groups:
                        parse_shape(
                            child,
                            z_index=child_idx,
                            level=level + 1,
                            parent_id=path,
                            path=child_path,
                        )

        for z_index, shape in enumerate(slide.shapes, start=1):
            parse_shape(shape, z_index=z_index, level=0, parent_id=None, path=f"{slide_index}.{z_index}")

        if image_info is None:
            slide_warnings.append("Slide image not available, pixel coordinates are null.")

        slide_result = SlideResult(
            slide_index=slide_index,
            slide_id=getattr(slide, "slide_id", None),
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            image=(
                None
                if image_info is None
                else {
                    "width_px": image_info.width_px,
                    "height_px": image_info.height_px,
                    "image_base64": image_info.base64_data,
                }
            ),
            elements=slide_elements,
            warnings=slide_warnings,
        )
        slides.append(slide_result)

    return slides, total_elements
